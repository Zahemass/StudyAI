import os
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
import filetype


class DocumentProcessor:
    def __init__(self):
        print("✅ DocumentProcessor initialized")
        self.supported_formats = [
            'pdf', 'docx', 'doc', 'pptx', 'ppt', 
            'txt', 'xlsx', 'xls', 'md', 'csv'
        ]

    def detect_file_type(self, file_path):
        """Detect file type"""
        kind = filetype.guess(file_path)
        
        if kind is None:
            # Fallback to extension
            ext = os.path.splitext(file_path)[1].lower().replace('.', '')
            return ext
        
        return kind.extension

    def extract_text(self, file_path):
        """Universal text extraction"""
        print(f"📄 Processing document: {file_path}")
        
        if not os.path.exists(file_path):
            raise Exception(f"File not found: {file_path}")
        
        # Detect file type
        file_type = self.detect_file_type(file_path)
        print(f"   Detected type: {file_type}")
        
        # Route to appropriate extractor
        if file_type == 'pdf':
            return self.extract_from_pdf(file_path)
        elif file_type in ['docx', 'doc']: 
            return self.extract_from_docx(file_path)
        elif file_type in ['pptx', 'ppt']:
            return self.extract_from_pptx(file_path)
        elif file_type == 'txt': 
            return self.extract_from_txt(file_path)
        elif file_type in ['xlsx', 'xls']: 
            return self.extract_from_xlsx(file_path)
        elif file_type == 'md': 
            return self.extract_from_txt(file_path)  # Markdown is just text
        elif file_type == 'csv': 
            return self.extract_from_csv(file_path)
        else:
            raise Exception(f"Unsupported file type: {file_type}")

    def extract_from_pdf(self, file_path):
        """Extract text from PDF"""
        print("   📕 Extracting from PDF...")
        text = ""
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n\n"
        except Exception as e:
            raise Exception(f"PDF extraction failed:  {str(e)}")
        
        if not text.strip():
            raise Exception("No text found in PDF")
        
        print(f"   ✅ Extracted {len(text)} characters from PDF")
        return text.strip()

    def extract_from_docx(self, file_path):
        """Extract text from Word document"""
        print("   📘 Extracting from DOCX...")
        text = ""
        
        try:
            doc = Document(file_path)
            
            # Extract from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n\n"
            
            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    text += row_text + "\n"
                text += "\n"
        
        except Exception as e:
            raise Exception(f"DOCX extraction failed: {str(e)}")
        
        if not text.strip():
            raise Exception("No text found in DOCX")
        
        print(f"   ✅ Extracted {len(text)} characters from DOCX")
        return text.strip()

    def extract_from_pptx(self, file_path):
        """Extract text from PowerPoint"""
        print("   📙 Extracting from PPTX...")
        text = ""
        
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"--- Slide {slide_num} ---\n\n"
                
                # Extract from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                    
                    # Extract from tables in slides
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join([cell.text for cell in row.cells])
                            text += row_text + "\n"
                
                text += "\n\n"
        
        except Exception as e:
            raise Exception(f"PPTX extraction failed: {str(e)}")
        
        if not text.strip():
            raise Exception("No text found in PPTX")
        
        print(f"   ✅ Extracted {len(text)} characters from PPTX")
        return text.strip()

    def extract_from_txt(self, file_path):
        """Extract text from TXT/MD files"""
        print("   📄 Extracting from TXT...")
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                text = file.read()
        except Exception as e:
            raise Exception(f"TXT extraction failed:  {str(e)}")
        
        if not text.strip():
            raise Exception("No text found in TXT")
        
        print(f"   ✅ Extracted {len(text)} characters from TXT")
        return text.strip()

    def extract_from_xlsx(self, file_path):
        """Extract text from Excel"""
        print("   📗 Extracting from XLSX...")
        text = ""
        
        try: 
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"--- Sheet: {sheet_name} ---\n\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
                
                text += "\n\n"
        
        except Exception as e:
            raise Exception(f"XLSX extraction failed: {str(e)}")
        
        if not text.strip():
            raise Exception("No text found in XLSX")
        
        print(f"   ✅ Extracted {len(text)} characters from XLSX")
        return text.strip()

    def extract_from_csv(self, file_path):
        """Extract text from CSV"""
        print("   📊 Extracting from CSV...")
        text = ""
        
        try:
            import csv
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                reader = csv.reader(file)
                for row in reader: 
                    text += " | ".join(row) + "\n"
        
        except Exception as e:
            raise Exception(f"CSV extraction failed: {str(e)}")
        
        if not text.strip():
            raise Exception("No text found in CSV")
        
        print(f"   ✅ Extracted {len(text)} characters from CSV")
        return text.strip()

    def get_supported_formats(self):
        """Return list of supported formats"""
        return self.supported_formats